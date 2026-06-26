#!/usr/bin/env python3
"""Gerar slides .pptx para apresentação do FilaClara."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SLIDES_DIR = "/home/deivi/Projetos/Metodologia-EAD/disciplinas/gcetens843-projeto-algoritmo-i/atividades/trabalhos/projeto-algoritmos"
OUTPUT = os.path.join(SLIDES_DIR, "apresentacao-filaclara.pptx")

# Cores do FilaClara
COR_FUNDO = RGBColor(0x0F, 0x17, 0x2A)  # azul escuro
COR_BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
COR_ACENTO = RGBColor(0x4F, 0x46, 0xE5)  # indigo
COR_VERDE = RGBColor(0x10, 0xB9, 0x81)
PR = Presentation()
PR.slide_width = Inches(13.333)
PR.slide_height = Inches(7.5)

def add_slide(title, content, accent=COR_BRANCO):
    slide = PR.slides.add_slide(PR.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COR_FUNDO
    
    # Titulo
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = accent
    
    # Conteudo
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(content):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(24)
        p2.font.color.rgb = COR_BRANCO
        p2.space_after = Pt(12)
    return slide

# --- Slides ---

# 1. Capa
slide = PR.slides.add_slide(PR.slide_layouts[6])
bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = COR_FUNDO
shapes = slide.shapes
# Titulo principal
txBox = shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(2))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "FilaClara"; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = COR_VERDE
p2 = tf.add_paragraph(); p2.text = "Filas Inteligentes para Serviços de Saúde com Acessibilidade"
p2.font.size = Pt(28); p2.font.color.rgb = COR_BRANCO; p2.space_before = Pt(12)
# Info
txBox2 = shapes.add_textbox(Inches(1.5), Inches(5), Inches(10), Inches(2))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p3 = tf2.paragraphs[0]; p3.text = "GCETENS843 — Projeto Algoritmo I"
p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
p4 = tf2.add_paragraph(); p4.text = "UFRB · CETENS · Bacharelado em Sistemas de Informação (EAD)"
p4.font.size = Pt(18); p4.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
p5 = tf2.add_paragraph(); p5.text = "Apresentação em 29/06/2026"
p5.font.size = Pt(18); p5.font.color.rgb = COR_ACENTO

# 2. Problema
add_slide("O Problema", [
    "• Filas presenciais em unidades de saúde geram desconforto,",
    "  ansiedade e perda de tempo para pacientes",
    "",
    "• Falta de transparência: paciente não sabe sua posição",
    "  nem o tempo estimado de espera",
    "",
    "• Barreiras de acessibilidade para PcD, TEA, idosos,",
    "  gestantes e pessoas com deficiência auditiva/visual",
    "",
    "• Soluções existentes são genéricas e não priorizam",
    "  acessibilidade desde o início do design",
], COR_VERDE)

# 3. Solução
add_slide("Nossa Solução: FilaClara", [
    "• Sistema de acompanhamento de filas em tempo real",
    "",
    "• Retirada de senha virtual pelo celular — sem necessidade",
    "  de ficar parado na recepção",
    "",
    "• Notificação antecipada quando o atendimento está próximo",
    "",
    "• Painel público transparente: mostra a sequência de chamadas",
    "  e reduz suspeita de favorecimento",
    "",
    "• Acessibilidade integrada: alto contraste, texto ampliado,",
    "  alerta visual, modo simples (TEA), leitura por voz",
])

# 4. Fluxograma
add_slide("Fluxograma do Sistema", [
    "1. Usuário entra no sistema (cadastro/login)",
    "2. Solicita o atendimento desejado",
    "3. É prioridade? → Fila Prioritária ou Convencional",
    "4. Sistema gera senha e posiciona na fila",
    "5. Topo da fila é chamado para atendimento",
    "6. 2º na fila recebe notificação para comparecer",
    "7. Sistema calcula e exibe tempo médio",
    "8. Usuário é atendido e excluído da fila",
    "9. Fila atualiza automaticamente",
], COR_VERDE)

# 5. Protótipo
add_slide("Protótipo de Telas", [
    "• Protótipo interativo disponível em:",
    "  github.com/Deivisan/Metodologia-EAD",
    "",
    "• Fluxo principal:",
    "  → Splash → Onboarding → Home",
    "  → Escolher clínica → Tipo de atendimento",
    "  → Senha virtual → Acompanhamento da fila",
    "",
    "• Funcionalidades de acessibilidade:",
    "  Alto contraste, texto ampliado, alerta visual,",
    "  vibração, modo simples (TEA), leitura por voz",
])

# 6. Acessibilidade
add_slide("Diferencial: Acessibilidade", [
    "♿ Acessibilidade não é um extra — é parte do design",
    "",
    "• Texto ampliado para baixa visão",
    "• Alto contraste para maior legibilidade",
    "• Modo simples para reduzir estímulos (TEA)",
    "• Alerta visual forte (não depende só de som)",
    "• Vibração para chamadas próximas",
    "• Leitura por voz para avisos principais",
    "",
    "▶ Aplicação real: Ausiane (TEA + deficiência auditiva)",
    "  trabalhou no NAI e trouxe essa perspectiva ao projeto",
], COR_VERDE)

# 7. Equipe
add_slide("Equipe", [
    "🧠 Deivison — Proposta + Documentação",
    "🏥 Rios — Pesquisa + Proposta (experiência em saúde)",
    "♿ Ausiane — Casos + Protótipo (acessibilidade e design)",
    "📚 Nubia — Pesquisa + Casos",
    "📊 Artur — Fluxograma",
    "🆕 Marcos Vinicius — Fluxograma",
    "🎤 Uélinton — Protótipo + Slides",
    "👤 Wallace — Slides",
], COR_ACENTO)

# 8. Obrigado
slide = PR.slides.add_slide(PR.slide_layouts[6])
bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = COR_FUNDO
txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Obrigado!"; p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = COR_VERDE
p2 = tf.add_paragraph(); p2.text = "Perguntas?"; p2.font.size = Pt(36); p2.font.color.rgb = COR_BRANCO; p2.space_before = Pt(20)
p3 = tf.add_paragraph(); p3.text = "github.com/Deivisan/Metodologia-EAD"; p3.font.size = Pt(20); p3.font.color.rgb = COR_ACENTO; p3.space_before = Pt(40)

PR.save(OUTPUT)
print(f"✅ Slides salvos em: {OUTPUT}")
